import os
import io
import re
import json
import hashlib
from datetime import datetime
from PIL import Image
import pytesseract
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from pdfminer.high_level import extract_text_to_fp
from .google_drive import GoogleDriveService
import asyncio
import logging

# Set up logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# Optional: Google API modules
try:
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaIoBaseDownload
    HAS_GOOGLE_API = True
except ImportError:
    HAS_GOOGLE_API = False

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

file_type_map = {
    'documents': ['docx', 'txt', 'doc', 'rtf', 'odt', 'pages', 'md', 'gdoc'],
    'spreadsheets': ['xlsx', 'xls', 'csv', 'ods', 'numbers', 'gsheet'],
    'presentations': ['pptx', 'ppt', 'odp', 'key', 'gslides'],
    'pdfs': ['pdf'],
    'images': ['jpg', 'jpeg', 'png', 'webp', 'gif', 'bmp', 'tiff', 'heic', 'gdraw'],
    'videos': ['mp4', 'mov', 'avi', 'wmv', 'flv', 'mkv', 'webm'],
    'audio': ['mp3', 'wav', 'ogg', 'm4a', 'wma'],
    'archives': ['zip', 'rar', '7z', 'tar', 'gz'],
    'code': ['py', 'js', 'java', 'cpp', 'h', 'cs', 'php', 'rb', 'swift', 'gs']
}

mime_type_map = {
    # Google Workspace types
    'application/vnd.google-apps.document': 'gdoc',
    'application/vnd.google-apps.spreadsheet': 'gsheet',
    'application/vnd.google-apps.presentation': 'gslides',
    'application/vnd.google-apps.drawing': 'gdraw',
    'application/vnd.google-apps.form': 'gform',
    'application/vnd.google-apps.script': 'gs',
    'application/vnd.google-apps.folder': 'folder',
    
    # Common document types
    'application/pdf': 'pdf',
    'application/msword': 'doc',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.ms-excel': 'xls',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'application/vnd.ms-powerpoint': 'ppt',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.oasis.opendocument.text': 'odt',
    'application/vnd.oasis.opendocument.spreadsheet': 'ods',
    'application/vnd.oasis.opendocument.presentation': 'odp',
    'application/x-iwork-pages-sffpages': 'pages',
    'application/x-iwork-numbers-sffnumbers': 'numbers',
    'application/x-iwork-keynote-sffkey': 'key',
    'text/markdown': 'md',
    'text/plain': 'txt',
    'text/rtf': 'rtf',
    
    # Image types
    'image/jpeg': 'jpg',
    'image/png': 'png',
    'image/gif': 'gif',
    'image/webp': 'webp',
    'image/bmp': 'bmp',
    'image/tiff': 'tiff',
    'image/heic': 'heic',
    
    # Video types
    'video/mp4': 'mp4',
    'video/quicktime': 'mov',
    'video/x-msvideo': 'avi',
    'video/x-ms-wmv': 'wmv',
    'video/webm': 'webm',
    'video/x-matroska': 'mkv',
    
    # Audio types
    'audio/mpeg': 'mp3',
    'audio/wav': 'wav',
    'audio/ogg': 'ogg',
    'audio/mp4': 'm4a',
    'audio/x-ms-wma': 'wma',
    
    # Archive types
    'application/zip': 'zip',
    'application/x-rar-compressed': 'rar',
    'application/x-7z-compressed': '7z',
    'application/x-tar': 'tar',
    'application/gzip': 'gz',
    
    # Text and code
    'text/javascript': 'js',
    'text/x-python': 'py',
    'text/x-java': 'java',
    'text/x-c': 'c',
    'text/x-cpp': 'cpp',
    'text/x-csharp': 'cs',
    'text/x-php': 'php',
    'text/x-ruby': 'rb',
    'text/x-swift': 'swift'
}

sensitive_keywords = {
    "pii": [
        "dob", "email", "phone", "address", "ssn", "personal", "pii", 
        "hipaa", "gdpr", "personally identifiable", "customer data",
        "personnel", "employee", "patient", "healthcare"
    ],
    "financial": [
        "credit", "bank", "amount", "revenue", "budget", "roi", "cost",
        "financial", "invoice", "payment", "expense", "profit", "pricing",
        "salary", "investment", "tax"
    ],
    "legal": [
        "license", "contract", "agreement", "legal", "compliance",
        "regulatory", "counsel", "policy", "policies", "terms",
        "regulation", "gdpr", "ccpa", "hipaa", "certification",
        "audit", "liability"
    ],
    "confidential": [
        "confidential", "internal use", "do not distribute", "sensitive",
        "security", "restricted", "proprietary", "classified", "private",
        "secret", "nda", "non-disclosure", "intellectual property",
        "trade secret", "internal only"
    ]
}

patterns = {
    # Matches common credit card formats (Visa, MC, Amex, Discover)
    "credit_card": r"(?:(?:4[0-9]{12}(?:[0-9]{3})?)|(?:5[1-5][0-9]{14})|(?:3[47][0-9]{13})|(?:6(?:011|5[0-9]{2})[0-9]{12}))",
    
    # Matches MM/YY or MM/YYYY with validation
    "expiry_date": r"(?:0[1-9]|1[0-2])\/(?:2[3-9]|[3-9][0-9])",
    
    # Matches SSN with required dashes and surrounding context
    "ssn": r"(?:SSN|Social Security)(?:[^0-9-])*\d{3}-\d{2}-\d{4}",
    
    # Matches email with common domains and validation
    "email": r"(?:[a-zA-Z0-9._%+-]+@(?:[a-zA-Z0-9-]+\.)+[a-zA-Z]{2,})",
    
    # Matches phone with required context and common formats
    "phone": r"(?<![\w/.:-])(?:(?:Phone|Tel|Mobile|Contact|Call|Fax)(?:[^0-9(])+)?(?:\+?1[-. ])?\(?[2-9][0-9]{2}\)?[-. ]?[2-9][0-9]{2}[-. ]?[0-9]{4}(?:\s*(?:ext|x)\.?\s*\d{1,5})?(?![-\d./@])",

    # Matches PA driver's license with validation
    "drivers_license": r"(?:Driver'?s? License|DL|License Number|License #)(?:[^0-9])*(?:[A-Z][0-9]{7}|[A-Z][0-9]{8}|[A-Z][0-9]{12}|\d{7,9}|[A-Z]\d{2}[-\s]?\d{3}[-\s]?\d{3}|[A-Z]\d{3}[-\s]?\d{3}[-\s]?\d{3}|[A-Z]{1,2}\d{4,7})",    
    
    # Matches address with validation and context
    "address_like": r"(?:Address|Location|Street)(?:[^0-9])*\d{1,5}\s[\w\s.]+(?:Street|St|Avenue|Ave|Road|Rd|Boulevard|Blvd|Lane|Ln|Drive|Dr|Circle|Cir|Court|Ct|Way|Place|Pl|Square|Sq)\b"
}

now = datetime.now()

def classify_by_age(modified_time):
    age_days = (now - modified_time).days
    if age_days <= 365:
        return "lessThanOneYear"
    elif age_days <= 1095:
        return "oneToThreeYears"
    else:
        return "moreThanThreeYears"

def initialize_structure():
    """Initialize the structure for file categorization."""
    return {
        "total_documents": 0,
        "total_sensitive": 0,
        "total_duplicates": 0,
        "file_types": {k: [] for k in file_type_map.keys() | {"others"}},
        "sensitive_info": {k: [] for k in sensitive_keywords.keys()},
        "duplicate_files": []
    }

def scan_text(text):
    """
    Scan text for sensitive information using keywords and patterns.
    Returns a dictionary of findings only if sensitive content is detected.
    """
    findings = {cat: [] for cat in sensitive_keywords}
    text_lower = text.lower()
    
    # Check for keywords in each category
    for cat, keywords in sensitive_keywords.items():
        for keyword in keywords:
            # Look for whole word matches only
            if re.search(r'\b' + re.escape(keyword.lower()) + r'\b', text_lower):
                findings[cat].append(keyword)
    
    # Check for pattern matches
    for label, pattern in patterns.items():
        if re.search(pattern, text):
            findings["pii"].append(label)
    
    # Only return categories that have findings
    return {k: v for k, v in findings.items() if v}

def extract_text_from_file(stream, file_type):
    try:
        if file_type == 'docx':
            doc = Document(stream)
            return "\n".join([p.text for p in doc.paragraphs])
        elif file_type == 'pptx':
            prs = Presentation(stream)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file_type in ['xlsx', 'xls']:
            wb = load_workbook(stream, read_only=True, data_only=True)
            return "\n".join([str(cell.value) for sheet in wb.worksheets for row in sheet.iter_rows() for cell in row if cell.value])
        elif file_type == 'pdf':
            output = io.StringIO()
            extract_text_to_fp(stream, output)
            return output.getvalue()
        elif file_type in ['jpg', 'jpeg', 'png', 'webp']:
            img = Image.open(stream)
            return pytesseract.image_to_string(img)
        elif file_type == 'txt':
            return stream.read().decode()
    except:
        return ""
    return ""

async def scan_files(source='local', path_or_drive_id='.', output_json='scan_report.json'):
    results = {
        "moreThanThreeYears": initialize_structure(),
        "oneToThreeYears": initialize_structure(),
        "lessThanOneYear": initialize_structure(),
        "scan_complete": False,
        "processed_files": 0,
        "total_files": 0,
        "total_duplicates": 0,
        "total_sensitive_files": 0,
        "failed_files": [],
        "content_hashes": {}
    }

    # Add logging for file type mapping
    logger.info(f"Using file type mapping: {file_type_map}")
    type_counts = {k: 0 for k in file_type_map.keys() | {"others"}}

    if source == 'local':
        all_files = []
        for root, _, files in os.walk(path_or_drive_id):
            for f in files:
                all_files.append(os.path.join(root, f))
        results["total_files"] = len(all_files)

        for filepath in all_files:
            try:
                ext = filepath.split('.')[-1].lower()
                modified_time = datetime.fromtimestamp(os.path.getmtime(filepath))
                age_group = classify_by_age(modified_time)
                file_type = next((k for k, v in file_type_map.items() if ext in v), "others")
                with open(filepath, 'rb') as f:
                    content = extract_text_from_file(f, ext)
                if not content:
                    continue
                results[age_group]["total_documents"] += 1
                results[age_group]["file_types"][file_type].append(filepath)
                findings = scan_text(content)
                if any(findings.values()):
                    results[age_group]["total_sensitive"] += 1
                    results["total_sensitive_files"] += 1
                    for k, v in findings.items():
                        results[age_group]["sensitive_info"][k].extend(v)
                results["processed_files"] += 1
            except:
                results["failed_files"].append(filepath)

    elif source == 'gdrive' and HAS_GOOGLE_API:
        drive_service = GoogleDriveService()
        if not drive_service.is_authenticated():
            raise ValueError("Not authenticated with Google Drive")

        try:
            files = await drive_service.list_directory(path_or_drive_id, recursive=True)
            results["total_files"] = len(files)
            logger.info(f"*** Total files found: {len(files)}")

            for file in files:
                try:
                    file_id = file['id']
                    name = file['name']
                    mime_type = file['mimeType']
                    
                    # Log file type categorization
                    logger.info(f"Processing file: {name} (mime_type: {mime_type})")
                    
                    # Get file extension from mime type or name
                    ext = mime_type_map.get(mime_type, None)
                    if not ext and '.' in name:
                        ext = name.split('.')[-1].lower()
                    
                    if not ext:
                        ext = 'others'

                    modified_time = datetime.fromisoformat(file['modifiedTime'].rstrip("Z"))
                    age_group = classify_by_age(modified_time)

                    # Determine file type category
                    file_type = 'others'
                    for category, extensions in file_type_map.items():
                        if ext in extensions:
                            file_type = category
                            break
                    
                    # Update type counts
                    type_counts[file_type] += 1
                    
                    # Add file to appropriate category
                    results[age_group]["total_documents"] += 1
                    results[age_group]["file_types"][file_type].append({
                        "id": file_id,
                        "name": name,
                        "mimeType": mime_type,
                        "modifiedTime": file['modifiedTime']
                    })

                    # Only scan content for text-based files
                    if file_type in ['documents', 'spreadsheets', 'presentations', 'pdfs']:
                        try:
                            content = await drive_service.get_file_content(file_id)
                            if content:
                                findings = scan_text(content)
                                if findings:
                                    results[age_group]["total_sensitive"] += 1
                                    results["total_sensitive_files"] += 1
                                    for k, v in findings.items():
                                        if v:  # Only add if there are findings
                                            results[age_group]["sensitive_info"][k].append({
                                                "file": {
                                                    "id": file_id,
                                                    "name": name,
                                                    "mimeType": mime_type,
                                                    "modifiedTime": file['modifiedTime']
                                                },
                                                "confidence": 0.8,
                                                "explanation": f"Found {', '.join(v)}",
                                                "categories": v
                                            })
                        except Exception as e:
                            logger.error(f"Error processing file content {name}: {str(e)}")
                    
                    results["processed_files"] += 1

                except Exception as e:
                    logger.error(f"Error with file {file.get('name', 'unknown')}: {str(e)}")
                    results["failed_files"].append({
                        "name": file.get('name', 'unknown'),
                        "error": str(e)
                    })

        except Exception as e:
            raise ValueError(f"Error accessing Google Drive: {str(e)}")

    # Log final type counts
    logger.info("File type counts:")
    for file_type, count in type_counts.items():
        logger.info(f"{file_type}: {count}")

    results["scan_complete"] = True
    with open(output_json, "w") as f:
        json.dump(results, f, indent=2)
    return results
